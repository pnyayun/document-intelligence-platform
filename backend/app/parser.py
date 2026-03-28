import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup

ALLOWED_EXTENSIONS = {
    'pdf', 'docx', 'txt', 'pptx', 'xlsx',
    'md', 'rtf', 'html', 'htm', 'csv'
}

def allowed_file(filename):
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_extension(filename):
    return filename.rsplit('.', 1)[1].lower()

def extract_text(filepath, filename):
    ext = get_extension(filename)
    extractors = {
        'pdf':  extract_pdf,
        'docx': extract_docx,
        'txt':  extract_txt,
        'md':   extract_txt,
        'csv':  extract_txt,
        'rtf':  extract_rtf,
        'pptx': extract_pptx,
        'xlsx': extract_xlsx,
        'html': extract_html,
        'htm':  extract_html,
    }
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f'Unsupported file type: .{ext}')
    return extractor(filepath)

def extract_pdf(filepath):
    text = ''
    page_count = 0
    with fitz.open(filepath) as doc:
        page_count = len(doc)
        for page in doc:
            text += page.get_text()
    return text.strip(), page_count

def extract_docx(filepath):
    doc = Document(filepath)
    text = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join([cell.text.strip() for cell in row.cells])
            if row_text.strip():
                text += '\n' + row_text
    return text.strip(), None

def extract_txt(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    return text.strip(), None

def extract_rtf(filepath):
    try:
        from striprtf.striprtf import rtf_to_text
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            rtf_content = f.read()
        text = rtf_to_text(rtf_content)
        return text.strip(), None
    except ImportError:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return text.strip(), None

def extract_pptx(filepath):
    prs = Presentation(filepath)
    text = ''
    slide_count = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        text += f'\n--- Slide {i+1} ---\n'
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text += shape.text.strip() + '\n'
    return text.strip(), slide_count

def extract_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    text = ''
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text += f'\n--- Sheet: {sheet} ---\n'
        for row in ws.iter_rows(values_only=True):
            row_text = ' | '.join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                text += row_text + '\n'
    wb.close()
    return text.strip(), None

def extract_html(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    soup = BeautifulSoup(content, 'html.parser')
    for tag in soup(['script', 'style', 'nav', 'footer']):
        tag.decompose()
    text = soup.get_text(separator='\n')
    text = '\n'.join([line.strip() for line in text.splitlines() if line.strip()])
    return text, None